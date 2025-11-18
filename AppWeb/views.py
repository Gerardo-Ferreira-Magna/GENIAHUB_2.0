
# --- Python estándar ---
import json
import io


# --- Django Core / HTTP / Utilidades ---
from django.http import JsonResponse
from django.shortcuts import render, redirect, get_object_or_404
from django.core.paginator import Paginator
from django.middleware.csrf import get_token
from django.utils.timezone import now, timedelta
from django.db.models import Q, Count, Avg


# --- Django Autenticación / Permisos ---
from django.contrib.auth import authenticate, login as auth_login, logout, get_user_model
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib.admin.views.decorators import staff_member_required

# --- Django Mensajes y Validaciones ---
from django.contrib import messages
from django.views.decorators.http import require_POST, require_http_methods

# --- Django Configuración / Sitio / Email ---
from django.conf import settings
from django.core.mail import send_mail
from django.contrib.sites.shortcuts import get_current_site
from django.contrib.sites.models import Site
from django.urls import reverse
from django.utils.crypto import get_random_string
from django.utils.text import slugify

# --- Modelos y Formularios ---
from .models_audit import AuditLog
from .forms import RegistroForm, LoginForm, RegistroEmpresaForm
from django.shortcuts import render
from .forms import PerfilForm
from .forms import ProyectoForm
from .models import Usuario, Carrera, Sede, Proyecto, RegistroEmpresa
from .forms import SolicitudEmpresaForm
from .models import SolicitudEmpresa
from django.utils import timezone
from django.http import HttpResponseForbidden
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas
from django.http import HttpResponse
import openai
from AppWeb.asistente_ai import consultar_openai
from django.views.decorators.csrf import csrf_exempt
from .asistente_ai import consultar_openai

# ============================================================
# VISTAS PÚBLICAS BÁSICAS
# ============================================================

def index(request):
    """Página de inicio."""
    return render(request, "webs/index.html")

def nosotros(request):
    """Página 'Sobre nosotros'."""
    return render(request, "webs/nosotros.html")

def tareas(request):
    """Vista de tareas (en desarrollo)."""
    return render(request, "webs/tareas.html")


# ============================================================
# 🏢 REGISTRO DE EMPRESA SIN LOGIN
# ============================================================

def _enviar_correo_estado(request, registro: RegistroEmpresa):
    """
    Enviar correo con el enlace de seguimiento de solicitud de empresa.
    Maneja automáticamente el dominio local si no hay Site configurado.
    """
    try:
        dominio = get_current_site(request).domain
    except Site.DoesNotExist:
        dominio = "127.0.0.1:8000"

    link = f"http://{dominio}{reverse('estado_solicitud', args=[registro.uuid_seguimiento])}"
    asunto = "Seguimiento de tu solicitud - GENIAHUB"
    mensaje = (
        f"Hola,\n\n"
        f"Gracias por registrar a {registro.nombre_empresa} en GENIAHUB.\n"
        f"Puedes revisar el estado de tu solicitud aquí:\n{link}\n\n"
        f"Saludos,\nEquipo GENIAHUB"
    )

    send_mail(
        subject=asunto,
        message=mensaje,
        from_email=getattr(settings, "DEFAULT_FROM_EMAIL", "no-reply@geniahub.cl"),
        recipient_list=[registro.email_tracking],
        fail_silently=True,
    )


def registro_empresa(request):
    """
    Permite enviar solicitud sin login. 
    Crea o reutiliza un usuario EMP inactivo y guarda un registro con estado 'PEN'.
    """
    if request.method == 'POST':
        form = RegistroEmpresaForm(request.POST, request.FILES)
        if form.is_valid():
            data = form.cleaned_data
            email = data['correo_contacto']
            UserModel = get_user_model()

            # Si ya existe una solicitud pendiente, mostrar modal "en proceso"
            if RegistroEmpresa.objects.filter(email_tracking=email, estado="PEN").exists():
                messages.warning(request, "Ya tienes una solicitud en proceso de revisión.")
                return render(request, 'webs/registro_empresa.html', {'form': form})

            # Crear o reutilizar usuario EMP inactivo
            user, created = UserModel.objects.get_or_create(
                email=email,
                defaults={
                    'nombre': data['nombre_empresa'],
                    'apellido_paterno': 'Empresa',
                    'apellido_materno': 'GENIAHUB',
                    'rut': get_random_string(9),  # dummy
                    'rol': 'EMP',
                    'is_active': False,
                }
            )

            # Si ya existía, asegurar estado EMP e inactividad
            if not created:
                if getattr(user, 'rol', None) != 'EMP':
                    user.rol = 'EMP'
                user.is_active = False
                user.save(update_fields=['rol', 'is_active'])

            # Crear registro
            registro = form.save(commit=False)
            registro.created_by = user
            registro.updated_by = user
            registro.email_tracking = email
            registro.estado = "PEN"
            registro.save()

            # Enviar correo
            _enviar_correo_estado(request, registro)

            # Mensaje para mostrar modal de éxito
            messages.success(request, "Tu solicitud fue enviada correctamente.")
            return render(request, "webs/solicitud_enviada.html", {"correo": email})
        else:
            messages.error(request, "⚠️ Revisa los campos e intenta nuevamente.")
    else:
        form = RegistroEmpresaForm()

    return render(request, 'webs/registro_empresa.html', {'form': form})


def estado_solicitud(request, uuid):
    """Permite a la empresa revisar el estado actual de su solicitud."""
    solicitud = get_object_or_404(RegistroEmpresa, uuid_seguimiento=uuid)
    return render(request, 'webs/estado_solicitud.html', {'solicitud': solicitud})


def reenviar_estado(request, uuid):
    """Permite reenviar el enlace de seguimiento al correo registrado."""
    solicitud = get_object_or_404(RegistroEmpresa, uuid_seguimiento=uuid)
    _enviar_correo_estado(request, solicitud)
    messages.success(request, "📧 Enlace reenviado al correo registrado.")
    return redirect('estado_solicitud', uuid=uuid)


# ============================================================
# 🧩 PANEL DE SOLICITUDES (SOLO ADMINISTRADORES)
# ============================================================

@staff_member_required
def panel_solicitudes(request):
    """Lista todas las solicitudes de empresas para revisión interna."""
    solicitudes = RegistroEmpresa.objects.all().order_by('-created_at')
    return render(request, 'webs/panel_solicitudes.html', {'solicitudes': solicitudes})


@staff_member_required
def cambiar_estado_solicitud(request, pk, nuevo_estado):
    """Permite a un administrador aprobar o rechazar una solicitud."""
    solicitud = get_object_or_404(RegistroEmpresa, pk=pk)
    solicitud.estado = nuevo_estado
    solicitud.save(update_fields=['estado'])
    _enviar_correo_estado(request, solicitud)
    messages.success(request, f"✅ Estado cambiado a {solicitud.get_estado_display()}")
    return redirect('panel_solicitudes')


# ============================================================
# 📚 PROYECTOS (DEMOS / VISTA PÚBLICA)
# ============================================================

@login_required
def proyectos(request):

    # JOIN reales con select_related
    proyectos = (
        Proyecto.objects
        .select_related("autor", "carrera", "sede")
        .order_by("-updated_at")
    )

    proyectos_list = []

    # Mapeo de estados → colores e íconos
    estado_map = {
        "BOR": ("secondary", "bi-pencil"),
        "REV": ("info", "bi-search"),
        "APR": ("success", "bi-check2-circle"),
        "ACT": ("primary", "bi-lightning-charge"),
        "DET": ("warning", "bi-pause-circle"),
        "FIN": ("dark", "bi-flag-checkered"),
        "PUB": ("success", "bi-megaphone"),
        "ARC": ("secondary", "bi-archive"),
    }

    for p in proyectos:

        estado_color, estado_icono = estado_map.get(
            p.estado,
            ("secondary", "bi-question-circle")
        )

        proyectos_list.append({
            "id": p.id,
            "nombre": p.titulo,
            "descripcion": p.descripcion,
            "estado": p.get_estado_display(),
            "estado_color": estado_color,
            "estado_icono": estado_icono,
            "docente": f"{p.autor.nombre} {p.autor.apellido_paterno}",
            "categoria": p.carrera.nombre if p.carrera else "Sin categoría",
            "fecha_modificacion": p.updated_at.strftime("%d-%m-%Y"),
        })

    # Recomendaciones placeholder (funciona igual)
    recomendaciones_ia = [
        {
            "titulo": "Proyectos recomendados",
            "descripcion": "Sugerencias basadas en tus intereses.",
            "icono": "bi-stars"
        },
        {
            "titulo": "Docentes activos",
            "descripcion": "Docentes con proyectos recientes.",
            "icono": "bi-person-badge"
        },
    ]

    return render(request, "webs/proyectos.html", {
        "proyectos": proyectos_list,
        "recomendaciones_ia": recomendaciones_ia,
    })


# ============================================================
# 👤 REGISTRO / LOGIN / LOGOUT DE USUARIOS
# ============================================================

def registro(request):
    """Registro de nuevos usuarios."""
    if request.method == "POST":
        form = RegistroForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, "✅ Usuario creado con éxito. Ya puedes iniciar sesión.")
            return redirect("login")
        else:
            messages.error(request, "❌ Corrige los errores indicados abajo.")
    else:
        form = RegistroForm()
    return render(request, "webs/registro.html", {"form": form})


def login(request):
    """Inicio de sesión."""
    if request.method == 'POST':
        form = LoginForm(request.POST)
        if form.is_valid():
            email = form.cleaned_data['email']
            password = form.cleaned_data['password']

            user = authenticate(request, email=email, password=password)
            if user is not None:
                auth_login(request, user)
                messages.success(request, f"👋 Bienvenido {user.nombre} {user.apellido_paterno}")
                return redirect('proyectos')
            else:
                messages.error(request, "❌ Correo o contraseña incorrectos.")
    else:
        form = LoginForm()
    return render(request, 'webs/login.html', {'form': form})


def logout_view(request):
    """Cerrar sesión."""
    logout(request)
    messages.success(request, "👋 Has cerrado sesión correctamente.")
    return redirect('login')


# ============================================================
# ⚙️ PANEL DE ADMINISTRACIÓN DE USUARIOS
# ============================================================

from django.contrib.auth.decorators import login_required
from django.shortcuts import render
from .models import Proyecto


@login_required
def panel(request):
    """Panel de proyectos: 
       - Admin ve todos los proyectos
       - Docente ve solo los suyos
    """

    # === FILTRO POR ROL ===
    if request.user.rol == "ADM":
        proyectos = Proyecto.objects.all().select_related("autor").order_by("-created_at")
    else:
        proyectos = Proyecto.objects.filter(
            autor=request.user
        ).select_related("autor").order_by("-created_at")

    # === CONTADORES PARA TARJETAS ===
    total_proyectos = proyectos.count()
    activos = proyectos.filter(estado="ACT").count()
    revision = proyectos.filter(estado="REV").count()
    aprobados = proyectos.filter(estado="APR").count()
    finalizados = proyectos.filter(estado="FIN").count()
    detenidos = proyectos.filter(estado="DET").count()
    borradores = proyectos.filter(estado="BOR").count()

    # === CONTEXTO ===
    contexto = {
        "proyectos": proyectos,
        "usuario": request.user,

        # Tarjetas resumen
        "total_proyectos": total_proyectos,
        "activos": activos,
        "revision": revision,
        "aprobados": aprobados,
        "finalizados": finalizados,
        "detenidos": detenidos,
        "borradores": borradores,
    }

    return render(request, "webs/panel.html", contexto)



def staff_required(view_func):
    """Decorador auxiliar para vistas staff personalizadas."""
    return user_passes_test(lambda u: u.is_authenticated and u.is_staff, login_url='/login/')(view_func)


@staff_required
def usuarios_panel(request):
    """Panel de gestión de usuarios para staff."""
    qs = Usuario.objects.all().order_by(request.GET.get('sort', '-date_joined'))
    exclude_super = request.GET.get('exclude_super', '1')
    if exclude_super == '1':
        qs = qs.exclude(is_superuser=True)

    # Filtros dinámicos (sin cambios)
    q = request.GET.get('q')
    if q:
        qs = qs.filter(
            Q(nombre__icontains=q) |
            Q(apellido_paterno__icontains=q) |
            Q(apellido_materno__icontains=q)
        )
    rut = request.GET.get('rut')
    if rut:
        qs = qs.filter(rut__icontains=rut)
    email = request.GET.get('email')
    if email:
        qs = qs.filter(email__icontains=email)
    rol = request.GET.get('rol')
    if rol:
        qs = qs.filter(rol=rol)
    activo = request.GET.get('activo')
    if activo in ('0', '1'):
        qs = qs.filter(is_active=(activo == '1'))
    staff = request.GET.get('staff')
    if staff in ('0', '1'):
        qs = qs.filter(is_staff=(staff == '1'))

    # --- NUEVOS CONTEOS PARA TARJETAS ---
    total_usuarios = Usuario.objects.exclude(is_superuser=True).count()
    activos = Usuario.objects.filter(is_active=True, is_superuser=False).count()
    inactivos = Usuario.objects.filter(is_active=False, is_superuser=False).count()
    docentes = Usuario.objects.filter(rol='DOC', is_superuser=False).count()
    estudiantes = Usuario.objects.filter(rol='EST', is_superuser=False).count()
    empresas = Usuario.objects.filter(rol='EMP', is_superuser=False).count()

    paginator = Paginator(qs, 10)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    admin_user_change_url = f"admin:{Usuario._meta.app_label}_{Usuario._meta.model_name}_change"
    admin_user_add_url = f"admin:{Usuario._meta.app_label}_{Usuario._meta.model_name}_add"

    csrf_token = get_token(request)

    return render(request, 'webs/usuarios_panel.html', {
        'page_obj': page_obj,
        'admin_user_change_url': admin_user_change_url,
        'admin_user_add_url': admin_user_add_url,
        'csrf_token': csrf_token,
        # Conteos para las tarjetas
        'total_usuarios': total_usuarios,
        'activos': activos,
        'inactivos': inactivos,
        'docentes': docentes,
        'estudiantes': estudiantes,
        'empresas': empresas,
    })



# ============================================================
# 🧾 API USUARIOS (DETALLE / UPDATE / DELETE)
# ============================================================

from django.db.models import ProtectedError

@staff_required
@require_POST
def usuario_delete_api(request, pk):
    """API: elimina usuario de forma segura y controlada."""
    try:
        u = get_object_or_404(Usuario, pk=pk)

        if u.is_superuser:
            return JsonResponse({"ok": False, "error": "❌ No se puede eliminar un superusuario."}, status=403)

        try:
            u.delete()
            return JsonResponse({"ok": True, "message": f"✅ Usuario '{u.nombre} {u.apellido_paterno}' eliminado correctamente."})
        except ProtectedError as e:
            print(f"[ERROR] No se puede eliminar el usuario {u.id}: {e}")
            return JsonResponse({"ok": False, "error": "Este usuario tiene registros asociados y no puede eliminarse."}, status=400)

    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({"ok": False, "error": f"Error interno: {str(e)}"}, status=500)


@staff_required
def usuario_detail_api(request, pk):
    """API: devuelve detalle de un usuario para edición."""
    try:
        u = get_object_or_404(Usuario, pk=pk)
        data = {
            "id": u.id,
            "nombre": u.nombre,
            "apellido_paterno": u.apellido_paterno,
            "apellido_materno": u.apellido_materno,
            "rut": u.rut,
            "email": u.email,
            "rol": u.rol,
            "is_active": u.is_active,
            "is_staff": u.is_staff,
            "date_joined": u.date_joined.isoformat(),
        }
        return JsonResponse({"ok": True, "usuario": data})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({"ok": False, "error": f"Error interno: {str(e)}"}, status=500)
    

@staff_required
@require_http_methods(["POST"])
def usuario_update_api(request, pk):
    """API: actualiza datos de un usuario existente."""
    try:
        data = json.loads(request.body.decode('utf-8'))
    except Exception:
        return JsonResponse({"ok": False, "error": "JSON inválido"}, status=400)

    u = get_object_or_404(Usuario, pk=pk)

    # Campos editables
    campos = ['nombre', 'apellido_paterno', 'apellido_materno', 'rut', 'email', 'rol', 'is_active', 'is_staff']
    for campo in campos:
        if campo in data:
            setattr(u, campo, data[campo])

    # Validar contraseña si fue enviada
    pwd = data.get('password')
    pwd_confirm = data.get('password_confirm')
    if pwd or pwd_confirm:
        if pwd != pwd_confirm:
            return JsonResponse({"ok": False, "error": "Las contraseñas no coinciden"}, status=400)
        if len(pwd) < 6:
            return JsonResponse({"ok": False, "error": "La contraseña debe tener al menos 6 caracteres"}, status=400)
        u.set_password(pwd)

    u.save()
    return JsonResponse({"ok": True, "message": "✅ Usuario actualizado correctamente", "usuario_id": u.id})


@staff_required
@require_POST
def usuario_delete_api(request, pk):
    """API: elimina un usuario (seguro, con control de errores)."""
    try:
        u = get_object_or_404(Usuario, pk=pk)

        # Seguridad: evitar eliminar superusuarios
        if u.is_superuser:
            return JsonResponse({"ok": False, "error": "No se puede eliminar un superusuario."}, status=403)

        nombre = f"{u.nombre} {u.apellido_paterno}".strip()
        u.delete()
        return JsonResponse({"ok": True, "message": f"🗑️ Usuario '{nombre}' eliminado correctamente."})

    except Usuario.DoesNotExist:
        return JsonResponse({"ok": False, "error": "Usuario no encontrado."}, status=404)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({"ok": False, "error": f"Error interno: {str(e)}"}, status=500)



@staff_required
@require_POST
def usuario_delete_api(request, pk):
    """API: elimina usuario (seguro para no borrar superusuarios)."""
    u = get_object_or_404(Usuario, pk=pk)
    if u.is_superuser:
        return JsonResponse({"ok": False, "error": "No se puede eliminar un superusuario"}, status=403)
    u.delete()
    return JsonResponse({"ok": True, "message": "Usuario eliminado"})


# ============================================================
@require_POST
def verificar_solicitud(request):
    """
    Verifica si ya existe una solicitud 'PEN' para el correo o el NIF/CIF.
    Devuelve: { exists_email, exists_nif, exists_any }
    """
    correo = (request.POST.get('correo') or '').strip()
    nif_cif = (request.POST.get('nif_cif') or '').strip()

    q = RegistroEmpresa.objects.filter(estado="PEN")
    exists_email = bool(correo) and q.filter(email_tracking__iexact=correo).exists()
    exists_nif   = bool(nif_cif) and q.filter(nif_cif__iexact=nif_cif).exists()

    return JsonResponse({
        "exists_email": exists_email,
        "exists_nif": exists_nif,
        "exists_any": exists_email or exists_nif
    })


@require_POST
def crear_solicitud(request):
    """
    Crea la solicitud real en BD usando el mismo Form (soporta archivo).
    Devuelve: { ok: True } o { ok: False, errors: {...} }
    """
    form = RegistroEmpresaForm(request.POST, request.FILES)
    if not form.is_valid():
        # Retorna errores del form para mostrarlos en el front
        return JsonResponse({"ok": False, "errors": form.errors}, status=400)

    data = form.cleaned_data

    # Crear o reutilizar usuario EMP inactivo
    email = data.get('correo_contacto')
    UserModel = get_user_model()
    user, created = UserModel.objects.get_or_create(
        email=email,
        defaults={
            'nombre': data.get('nombre_empresa') or 'Empresa',
            'apellido_paterno': 'Empresa',
            'apellido_materno': 'GENIAHUB',
            'rut': get_random_string(9),  # dummy
            'rol': 'EMP',
            'is_active': False,
        }
    )
    if not created:
        changed = False
        if getattr(user, 'rol', None) != 'EMP':
            user.rol = 'EMP'
            changed = True
        if user.is_active:
            user.is_active = False
            changed = True
        if changed:
            user.save(update_fields=['rol', 'is_active'])

    # Crear el registro
    registro = form.save(commit=False)
    registro.created_by = user
    registro.updated_by = user
    registro.email_tracking = email
    registro.estado = "PEN"
    registro.save()

    # Si quieres: _enviar_correo_estado(request, registro)

    return JsonResponse({"ok": True})


# === PANEL DE SOLICITUDES ===
def panel_solicitudes(request):
    solicitudes = RegistroEmpresa.objects.all()

    # Conteos directos optimizados en SQL
    total = solicitudes.count()
    aprobadas = solicitudes.filter(estado='APR').count()
    pendientes = solicitudes.filter(estado='PEN').count()
    rechazadas = solicitudes.filter(estado='REJ').count()

    context = {
        'solicitudes': solicitudes,
        'total': total,
        'aprobadas': aprobadas,
        'pendientes': pendientes,
        'rechazadas': rechazadas,
    }
    return render(request, 'webs/panel_solicitudes.html', context)


# === EDITAR SOLICITUD ===
def editar_solicitud(request, pk):
    """Permite editar los datos de una solicitud (empresa, correo, NIF/CIF, estado, motivo)"""
    solicitud = get_object_or_404(RegistroEmpresa, pk=pk)

    if request.method == 'POST':
        solicitud.nombre_empresa = request.POST.get('nombre_empresa', solicitud.nombre_empresa)
        solicitud.nif_cif = request.POST.get('nif_cif', solicitud.nif_cif)
        solicitud.email_tracking = request.POST.get('correo_contacto', solicitud.email_tracking)
        solicitud.estado = request.POST.get('estado', solicitud.estado)
        solicitud.motivo_rechazo = request.POST.get('motivo_rechazo', solicitud.motivo_rechazo)

        solicitud.save()

        # Mensaje dinámico según el nuevo estado
        if solicitud.estado == 'APR':
            messages.success(request, f"✅ La solicitud de {solicitud.nombre_empresa} fue aprobada correctamente.")
        elif solicitud.estado == 'REJ':
            messages.warning(request, f"⚠️ La solicitud de {solicitud.nombre_empresa} fue rechazada.")
        else:
            messages.info(request, f"✏️ La solicitud de {solicitud.nombre_empresa} fue actualizada.")

        return redirect('panel_solicitudes')

    return redirect('panel_solicitudes')


# === CAMBIAR ESTADO (Aprobar / Rechazar) ===
def cambiar_estado_solicitud(request, pk, nuevo_estado):
    """Cambia el estado de una solicitud a Aprobado o Rechazado"""
    solicitud = get_object_or_404(RegistroEmpresa, pk=pk)

    # Validar el nuevo estado
    if nuevo_estado not in ['APR', 'REJ', 'PEN']:
        messages.error(request, "Estado inválido.")
        return redirect('panel_solicitudes')

    solicitud.estado = nuevo_estado

    # Si se rechaza sin motivo, asignar texto por defecto
    if nuevo_estado == 'REJ' and not solicitud.motivo_rechazo:
        solicitud.motivo_rechazo = "Rechazada por revisión interna."

    solicitud.save()

    if nuevo_estado == 'APR':
        messages.success(request, f"✅ Solicitud de {solicitud.nombre_empresa} aprobada.")
    elif nuevo_estado == 'REJ':
        messages.warning(request, f"⚠️ Solicitud de {solicitud.nombre_empresa} rechazada.")
    else:
        messages.info(request, f"📋 Estado actualizado para {solicitud.nombre_empresa}.")

    return redirect('panel_solicitudes')


@staff_required
@require_http_methods(["POST"])
def usuario_create_api(request):
    """API: crea un nuevo usuario (staff/admin)."""
    try:
        data = json.loads(request.body.decode('utf-8'))
    except Exception:
        return JsonResponse({"ok": False, "error": "JSON inválido"}, status=400)

    campos_obligatorios = ['nombre', 'apellido_paterno', 'rut', 'email', 'rol']
    if not all(data.get(c) for c in campos_obligatorios):
        return JsonResponse({"ok": False, "error": "Faltan campos obligatorios"}, status=400)

    # Validar contraseñas
    pwd = data.get('password')
    pwd_confirm = data.get('password_confirm')
    if not pwd or not pwd_confirm:
        return JsonResponse({"ok": False, "error": "Debe ingresar y confirmar la contraseña"}, status=400)
    if pwd != pwd_confirm:
        return JsonResponse({"ok": False, "error": "Las contraseñas no coinciden"}, status=400)
    if len(pwd) < 6:
        return JsonResponse({"ok": False, "error": "La contraseña debe tener al menos 6 caracteres"}, status=400)

    # Crear usuario
    User = get_user_model()
    if User.objects.filter(email=data['email']).exists():
        return JsonResponse({"ok": False, "error": "El correo ya está registrado"}, status=400)

    u = User(
        nombre=data['nombre'],
        apellido_paterno=data['apellido_paterno'],
        apellido_materno=data.get('apellido_materno', ''),
        rut=data['rut'],
        email=data['email'],
        rol=data['rol'],
        is_active=data.get('is_active', True),
        is_staff=data.get('is_staff', False),
    )
    u.set_password(pwd)
    u.save()

    return JsonResponse({"ok": True, "usuario_id": u.id, "message": "Usuario creado correctamente"})


@staff_required
@require_http_methods(["POST"])
def usuario_create_api(request):
    """API: crea un nuevo usuario (solo staff)."""
    try:
        data = json.loads(request.body.decode('utf-8'))
    except Exception:
        return JsonResponse({"ok": False, "error": "JSON inválido"}, status=400)

    campos_obligatorios = ['nombre', 'apellido_paterno', 'apellido_materno', 'rut', 'email', 'rol']
    for campo in campos_obligatorios:
        if not data.get(campo):
            return JsonResponse({"ok": False, "error": f"El campo '{campo}' es obligatorio."}, status=400)

    password = data.get('password')
    password_confirm = data.get('password_confirm')
    if not password or not password_confirm:
        return JsonResponse({"ok": False, "error": "Debe ingresar y confirmar una contraseña."}, status=400)
    if password != password_confirm:
        return JsonResponse({"ok": False, "error": "Las contraseñas no coinciden."}, status=400)
    if len(password) < 6:
        return JsonResponse({"ok": False, "error": "La contraseña debe tener al menos 6 caracteres."}, status=400)

    # Crear usuario
    UsuarioModel = get_user_model()
    if UsuarioModel.objects.filter(email=data['email']).exists():
        return JsonResponse({"ok": False, "error": "Ya existe un usuario con ese correo."}, status=400)

    u = UsuarioModel.objects.create(
        nombre=data['nombre'],
        apellido_paterno=data['apellido_paterno'],
        apellido_materno=data['apellido_materno'],
        rut=data['rut'],
        email=data['email'],
        rol=data['rol'],
        is_active=data.get('is_active', True),
        is_staff=data.get('is_staff', False),
    )
    u.set_password(password)
    u.save()

    return JsonResponse({"ok": True, "message": "Usuario creado correctamente", "usuario_id": u.id})


@login_required
def perfil_estudiante(request):
    user = (
        Usuario.objects
        .select_related("sede", "carrera")
        .get(pk=request.user.pk)
    )

    habilidades = user.habilidades.split(",") if user.habilidades else []
    industrias = user.industrias_interes.split(",") if user.industrias_interes else []
    tecnologias = user.tecnologias_preferidas.split(",") if user.tecnologias_preferidas else []

    return render(request, "webs/perfil_estudiante.html", {
        "usuario": user,
        "habilidades": [h.strip() for h in habilidades if h.strip()],
        "industrias": [i.strip() for i in industrias if i.strip()],
        "tecnologias": [t.strip() for t in tecnologias if t.strip()],
    })


@login_required
def editar_perfil(request):
    if request.method == "POST":
        user = request.user
        user.sobre_mi = request.POST.get("sobre_mi", "")
        user.habilidades = request.POST.get("habilidades", "")
        user.experiencia = request.POST.get("experiencia", "")
        user.industrias_interes = request.POST.get("industrias_interes", "")
        user.tecnologias_preferidas = request.POST.get("tecnologias_preferidas", "")
        user.save()
        return JsonResponse({"ok": True})
    return JsonResponse({"ok": False}, status=400)


@login_required
def busqueda_perfiles(request):
    """Vista principal para docentes (server-rendered, tarjetas oscuras)."""
    if request.user.rol != "DOC" and not request.user.is_superuser:
        messages.warning(request, "⚠️ Solo los docentes pueden acceder a esta vista.")
        return redirect("panel")

    # Filtros
    q = request.GET.get("q", "").strip()
    carrera_filtro = request.GET.get("carrera", "").strip()
    habilidad_filtro = request.GET.get("habilidad", "").strip()

    qs = (
        Usuario.objects.filter(rol="EST")
        .select_related("carrera", "sede")   # <-- ¡IMPORTANTE!
    )

    if q:
        qs = qs.filter(
            Q(nombre__icontains=q) |
            Q(apellido_paterno__icontains=q) |
            Q(apellido_materno__icontains=q) |
            Q(sobre_mi__icontains=q) |
            Q(habilidades__icontains=q) |
            Q(tecnologias_preferidas__icontains=q)
        )
    if carrera_filtro:
        qs = qs.filter(carrera__nombre__icontains=carrera_filtro)

    if habilidad_filtro:
        qs = qs.filter(habilidades__icontains=habilidad_filtro)

    qs = qs.order_by("apellido_paterno", "nombre")

    # Construcción de los datos
    estudiantes = []
    for e in qs:
        habs = [h.strip() for h in (e.habilidades or "").split(",") if h.strip()]
        inds = [i.strip() for i in (e.industrias_interes or "").split(",") if i.strip()]
        tech = [t.strip() for t in (e.tecnologias_preferidas or "").split(",") if t.strip()]

        estudiantes.append({
            "id": e.id,
            "nombre": f"{e.nombre} {e.apellido_paterno}",
            "email": e.email,
            "rut": e.rut,
            "carrera": e.carrera.nombre if e.carrera else "Sin carrera",
            "sede": e.sede.nombre if e.sede else "Sin sede",   # <<<<<< 🔥 AQUI EL FIX
            "sobre_mi": e.sobre_mi or "",
            "experiencia": e.experiencia or "",
            "habilidades": habs,
            "habilidades3": habs[:3],
            "industrias": inds,
            "tecnologias": tech,
        })

    # Para los filtros del formulario
    carreras = sorted({c.nombre for c in Carrera.objects.all()})
    habilidades = sorted({h for e in estudiantes for h in e["habilidades"]})

    return render(request, "webs/busqueda_perfiles.html", {
        "estudiantes": estudiantes,
        "carreras": carreras,
        "habilidades": habilidades,
        "query": q,
        "carrera_filtro": carrera_filtro,
        "habilidad_filtro": habilidad_filtro,
    })



@login_required
def buscar_perfiles_ajax(request):
    """Devuelve perfiles de estudiantes en JSON."""
    if request.user.rol not in ["DOC", "ADM"]:
        return JsonResponse({"error": "No autorizado"}, status=403)

    query = request.GET.get("q", "")
    carrera = request.GET.get("carrera", "")
    habilidad = request.GET.get("habilidad", "")

    estudiantes = Usuario.objects.filter(rol="EST")

    if query:
        estudiantes = estudiantes.filter(
            Q(nombre__icontains=query)
            | Q(apellido_paterno__icontains=query)
            | Q(habilidades__icontains=query)
            | Q(tecnologias_preferidas__icontains=query)
            | Q(sobre_mi__icontains=query)
        )
    if carrera:
        estudiantes = estudiantes.filter(carrera__nombre__icontains=carrera)
    if habilidad:
        estudiantes = estudiantes.filter(habilidades__icontains=habilidad)

    data = []
    for e in estudiantes:
        data.append({
            "id": e.id,
            "nombre": f"{e.nombre} {e.apellido_paterno}",
            "carrera": e.carrera.nombre if hasattr(e, 'carrera') and e.carrera else "Sin carrera",
            "habilidades": [h.strip() for h in (e.habilidades or "").split(",") if h.strip()],
            "email": e.email,
            "sobre_mi": e.sobre_mi or "",
            "experiencia": e.experiencia or "",
            "industrias": e.industrias_interes or "",
            "tecnologias": e.tecnologias_preferidas or "",
        })

    return JsonResponse({"resultados": data})


@login_required
def crear_proyecto(request):
    if request.user.rol not in ["DOC", "ADM"]:
        messages.error(request, "No tienes permisos para crear proyectos.")
        return redirect("proyectos")

    if request.method == "POST":
        form = ProyectoForm(request.POST, request.FILES)
        if form.is_valid():

            proyecto = form.save(commit=False)

            # Convertimos la fecha del form al AÑO
            fecha = form.cleaned_data["fecha_proyecto"]
            proyecto.anio = fecha.year

            # Auditoría
            proyecto.autor = request.user
            proyecto.created_by = request.user
            proyecto.updated_by = request.user

            proyecto.save()
            messages.success(request, "Proyecto creado exitosamente.")
            return redirect("proyectos")
        else:
            print(form.errors)   # Para debug

    else:
        form = ProyectoForm()

    context = {
        "form": form,
    }
    return render(request, "webs/creacion_proyecto.html", context)


@login_required
@require_POST
def proyecto_editar_modal(request):
    proyecto_id = request.POST.get("proyecto_id")
    proyecto = get_object_or_404(Proyecto, id=proyecto_id)

    # Permisos
    if request.user != proyecto.autor and request.user.rol != "ADM":
        return JsonResponse({"ok": False, "error": "No autorizado"}, status=403)

    proyecto.titulo = request.POST.get("titulo")
    proyecto.resumen = request.POST.get("resumen")
    proyecto.descripcion = request.POST.get("descripcion")
    proyecto.estado = request.POST.get("estado")

    # Regenerar slug
    proyecto.slug = slugify(f"{proyecto.titulo}-{proyecto.anio}")

    proyecto.save()

    return JsonResponse({
        "ok": True,
        "message": "Proyecto actualizado correctamente.",
        "titulo": proyecto.titulo,
        "estado": proyecto.estado,
        "resumen": proyecto.resumen,
        "descripcion": proyecto.descripcion
    })


@login_required
@require_POST
def proyecto_eliminar_modal(request):
    proyecto_id = request.POST.get("proyecto_id")
    proyecto = get_object_or_404(Proyecto, id=proyecto_id)

    # Permisos
    if request.user != proyecto.autor and request.user.rol != "ADM":
        return JsonResponse({"ok": False, "error": "No tienes permisos."}, status=403)

    proyecto.delete()

    return JsonResponse({
        "ok": True,
        "message": f"Proyecto '{proyecto.titulo}' eliminado correctamente."
    })


@login_required
def solicitud_empresa_crear(request):
    from django.utils import timezone

    if request.method == "POST":
        form = ProyectoForm(request.POST, request.FILES)

        if form.is_valid():
            proyecto = form.save(commit=False)

            # Datos fijos del flujo empresa
            proyecto.tipo = "EMP"
            proyecto.estado = "BOR"

            # Autor del proyecto
            proyecto.autor = request.user

            # Auditoría obligatoria
            proyecto.created_by = request.user
            proyecto.updated_by = request.user

            # Sede automática
            if request.user.sede_id:
                proyecto.sede = request.user.sede

            # Año desde la fecha
            fecha = form.cleaned_data.get("fecha_proyecto")
            if fecha:
                proyecto.anio = fecha.year

            proyecto.save()

            messages.success(request, "Solicitud creada correctamente.")
            return redirect("proyectos")

        else:
            print("ERRORES DEL FORM:", form.errors)

    else:
        form = ProyectoForm(initial={
            "tipo": "EMP",
            "estado": "BOR",
            "fecha_proyecto": timezone.now().date(),
            "sede": request.user.sede_id,
        })

    return render(request, "webs/solicitud_empresa.html", {"form": form})


@login_required
def empresa_estado_solicitud(request):

    # ADMIN ve todo
    if request.user.rol == "ADM":
        solicitudes = SolicitudEmpresa.objects.all()

    # EMPRESA ve solo las suyas
    elif request.user.rol == "EMP":
        solicitudes = SolicitudEmpresa.objects.filter(empresa=request.user)

    # OTROS no pueden ver
    else:
        solicitudes = SolicitudEmpresa.objects.none()

    solicitudes = solicitudes.select_related(
            "empresa",
            "asignacion",
            "asignacion__docente_responsable"
        ).prefetch_related(
            "asignacion__estudiantes"
        )

    contexto = {
        "solicitudes": solicitudes,
        "total": solicitudes.count(),
        "recibidas": solicitudes.filter(estado="REC").count(),
        "revision": solicitudes.filter(estado="REV").count(),
        "aprobadas": solicitudes.filter(estado="APR").count(),
        "asignadas": solicitudes.filter(estado="ASI").count(),
        "finalizadas": solicitudes.filter(estado="FIN").count(),
    }

    return render(request, "webs/empresa_estado_solicitud.html", contexto)


@login_required
def crear_solicitud_empresa(request):

    if request.user.rol != "EMP":
        messages.error(request, "Solo las empresas pueden enviar solicitudes.")
        return redirect("index")

    if request.method == "POST":
        form = SolicitudEmpresaForm(request.POST, request.FILES)

        if form.is_valid():
            solicitud = form.save(commit=False)

            solicitud.empresa = request.user
            solicitud.created_by = request.user
            solicitud.updated_by = request.user

            solicitud.save()

            messages.success(request, "Solicitud enviada correctamente.")
            return redirect("empresa_estado_solicitud")

        else:
            print("ERRORES FORM:", form.errors)

    else:
        form = SolicitudEmpresaForm()

    return render(request, "webs/solicitud_empresa.html", {"form": form})


@login_required
def proyectos_tabla(request):

    # =======================
    # SOLO proyectos empresariales
    # =======================
    proyectos = Proyecto.objects.select_related("autor").filter(tipo="EMP")

    # =======================
    # KPIs SOLO EMPRESA
    # =======================
    total_emp = proyectos.count()
    aprobados_emp = proyectos.filter(estado="APR").count()
    revision_emp = proyectos.filter(estado="REV").count()
    finalizados_emp = proyectos.filter(estado="FIN").count()
    detenidos_emp = proyectos.filter(estado="DET").count()

    contexto = {
        "proyectos": proyectos.order_by("-created_at"),

        # Estados para el filtro del select
        "proyecto_estados": Proyecto.ESTADO,

        # KPI SOLO EMPRESA
        "total_emp": total_emp,
        "aprobados_emp": aprobados_emp,
        "revision_emp": revision_emp,
        "finalizados_emp": finalizados_emp,
        "detenidos_emp": detenidos_emp,
    }

    return render(request, "webs/proyectos_tabla.html", contexto)


@login_required
def asistente_ia(request):
    filtro = request.GET.get("buscar", "")

    # usar Proyecto igual que en panel_proyectos
    qs = Proyecto.objects.select_related("autor").order_by("-created_at")

    if filtro:
        qs = qs.filter(
            Q(titulo__icontains=filtro) |
            Q(resumen__icontains=filtro) |
            Q(descripcion__icontains=filtro) |
            Q(autor__nombre__icontains=filtro) |
            Q(autor__apellido_paterno__icontains=filtro)
        )

    return render(request, "webs/asistente_ia.html", {
        "proyectos": qs,
        "buscar": filtro,
    })


@csrf_exempt
def asistente_ia_api(request):
    """API que recibe mensaje y responde con IA."""
    if request.method == "POST":
        data = json.loads(request.body)
        mensaje = data.get("mensaje", "")
        respuesta = consultar_openai(mensaje)
        return JsonResponse({"respuesta": respuesta})

    return JsonResponse({"error": "Método no permitido"}, status=405)
    

# ---- GENERAR PDF ----
def generar_pdf(request):
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer)
    p.drawString(100, 750, "Informe generado con IA - GENIAHUB")
    p.drawString(100, 720, "Proyecto: " + request.GET.get("titulo", "Sin título"))
    p.drawString(100, 700, "Resumen: " + request.GET.get("resumen", "Sin resumen"))
    p.save()

    buffer.seek(0)
    return HttpResponse(buffer, content_type='application/pdf')

# ---- GENERAR POWERPOINT ----
def generar_ppt(request):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Título
    slide = slide_layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = "Presentación IA - GENIAHUB"
    slide.placeholders[1].text = request.GET.get("resumen", "No hay contenido aún.")

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return HttpResponse(output, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

# ---- GENERAR EXCEL ----
def generar_excel(request):
    wb = Workbook()
    ws = wb.active
    ws.title = "Resumen IA"

    ws.append(["Proyecto", request.GET.get("titulo", "Sin título")])
    ws.append(["Descripción", request.GET.get("resumen", "No hay resumen")])
    ws.append(["Generado por", "Asistente IA GENIAHUB"])

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    return HttpResponse(
        output,
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )